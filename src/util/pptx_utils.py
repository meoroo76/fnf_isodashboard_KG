"""PowerPoint 유틸리티 — 슬라이드, 텍스트박스, KPI 카드, 테이블"""

from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── 색상 테마 상수 ───────────────────────────────────────────
# 기준: FNF_AI_Native_OptA_WhiteMinimal.pptx

THEME_DARK = {
    "bg": RGBColor(0x1E, 0x29, 0x3B),
    "primary": RGBColor(0x00, 0x6A, 0xE6),
    "accent": RGBColor(0xFF, 0xA0, 0x00),
    "card": RGBColor(0x2A, 0x36, 0x4A),
    "table_header": RGBColor(0x00, 0x6A, 0xE6),
    "table_row1": RGBColor(0x23, 0x2E, 0x3F),
    "table_row2": RGBColor(0x2A, 0x36, 0x4A),
    "text": RGBColor(0xFF, 0xFF, 0xFF),
    "text_sub": RGBColor(0x8A, 0x8A, 0x9A),
    "border": RGBColor(0x3A, 0x46, 0x5A),
}

THEME_LIGHT = {
    "bg": RGBColor(0xFA, 0xFA, 0xFC),
    "primary": RGBColor(0x00, 0x6A, 0xE6),
    "accent": RGBColor(0xFF, 0xA0, 0x00),
    "card": RGBColor(0xF0, 0xF2, 0xF5),
    "card_blue": RGBColor(0xE3, 0xF0, 0xFF),
    "card_green": RGBColor(0xE3, 0xF8, 0xF0),
    "card_orange": RGBColor(0xFF, 0xF3, 0xE0),
    "table_header": RGBColor(0x00, 0x6A, 0xE6),
    "table_row1": RGBColor(0xFF, 0xFF, 0xFF),
    "table_row2": RGBColor(0xF0, 0xF2, 0xF5),
    "text": RGBColor(0x2D, 0x2D, 0x2D),
    "text_sub": RGBColor(0x8A, 0x8A, 0x9A),
    "divider": RGBColor(0xC0, 0xC4, 0xCC),
    "divider_light": RGBColor(0xE8, 0xEA, 0xEE),
    "border": RGBColor(0xC0, 0xC4, 0xCC),
}

COLOR_POSITIVE = RGBColor(0x00, 0xB8, 0x94)
COLOR_NEGATIVE = RGBColor(0xFF, 0x63, 0x52)
COLOR_WARN = RGBColor(0xFF, 0xA0, 0x00)
COLOR_PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
COLOR_TEAL = RGBColor(0x00, 0x96, 0x88)

FONT_NAME = "Arial"

# ── 폰트 크기 기준 (pt) ──────────────────────────────────────
# 기준 디자인: FNF_AI_Native_OptA_WhiteMinimal.pptx
FONT_SIZES = {
    "cover_title": 38,
    "cover_subtitle": 18,
    "cover_date": 14,
    "title_bar": 20,
    "title_bar_sub": 12,
    "section_divider": 13,
    "kpi_label": 12,
    "kpi_value": 24,
    "kpi_sub": 11,
    "table_cell": 12,
    "bullet_title": 13,
    "bullet_item": 11,
}


# ── 프레젠테이션 생성 ────────────────────────────────────────

def create_presentation(
    width: float = 13.333, height: float = 7.5,
) -> Presentation:
    """16:9 와이드 프레젠테이션 생성."""
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    return prs


# ── 슬라이드 배경 ────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor) -> None:
    """슬라이드 배경색 설정."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── 텍스트 박스 ──────────────────────────────────────────────

def add_textbox(
    slide, left, top, width, height, text: str,
    font_size: int = 12, bold: bool = False,
    color: RGBColor = RGBColor(0x33, 0x33, 0x33),
    alignment=PP_ALIGN.LEFT,
    font_name: str = FONT_NAME,
    anchor=MSO_ANCHOR.TOP,
) -> Any:
    """텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


# ── 도형 ─────────────────────────────────────────────────────

def add_rounded_rect(
    slide, left, top, width, height,
    fill_color: RGBColor,
    border_color: RGBColor | None = None,
) -> Any:
    """둥근 사각형 추가."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_rect(
    slide, left, top, width, height,
    fill_color: RGBColor,
) -> Any:
    """사각형 추가 (내부용)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ── KPI 카드 ─────────────────────────────────────────────────

def add_kpi_card(
    slide, left, top, width, height,
    label: str, value: str, sub_text: str,
    accent_color: RGBColor,
    card_bg: RGBColor = THEME_LIGHT["card"],
    text_color: RGBColor = THEME_LIGHT["text"],
    label_color: RGBColor | None = None,
    label_fontsize: int | None = None,
    value_fontsize: int | None = None,
    sub_fontsize: int | None = None,
) -> Any:
    """KPI 카드 추가.

    상단 액센트 라인 + 라벨 + 값 + 보조 텍스트 구성.
    """
    label_color = label_color or RGBColor(0x8A, 0x8A, 0x9A)
    _lf = label_fontsize or FONT_SIZES["kpi_label"]
    _vf = value_fontsize or FONT_SIZES["kpi_value"]
    _sf = sub_fontsize or FONT_SIZES["kpi_sub"]
    card = add_rounded_rect(slide, left, top, width, height, card_bg)

    # 상단 액센트 라인
    add_rounded_rect(
        slide, left + Inches(0.15), top + Inches(0.1),
        Inches(0.5), Inches(0.04), accent_color,
    )
    # 라벨
    add_textbox(
        slide, left + Inches(0.15), top + Inches(0.2),
        width - Inches(0.3), Inches(0.3),
        label, font_size=_lf, color=label_color,
    )
    # 값
    add_textbox(
        slide, left + Inches(0.15), top + Inches(0.45),
        width - Inches(0.3), Inches(0.4),
        value, font_size=_vf, color=text_color, bold=True,
    )
    # 보조 텍스트
    add_textbox(
        slide, left + Inches(0.15), top + Inches(0.85),
        width - Inches(0.3), Inches(0.25),
        sub_text, font_size=_sf, color=accent_color, bold=True,
    )
    return card


# ── 테이블 ───────────────────────────────────────────────────

def create_styled_table(
    slide, left, top, width,
    data: list[list[str]],
    col_widths: list | None = None,
    theme: str = "dark",
    cell_fontsize: int | None = None,
) -> Any:
    """스타일링된 테이블 생성.

    Args:
        data: 2D 리스트. 첫 행은 헤더.
        col_widths: 각 컬럼 너비 (pptx 단위). None이면 균등 분배.
        theme: "dark" 또는 "light"
        cell_fontsize: 셀 폰트 크기 (None이면 FONT_SIZES 기본값)
    """
    t = THEME_DARK if theme == "dark" else THEME_LIGHT
    _cf = cell_fontsize or FONT_SIZES["table_cell"]
    rows = len(data)
    cols = len(data[0]) if data else 0
    if rows == 0 or cols == 0:
        return None

    row_h = Inches(0.35)
    table_shape = slide.shapes.add_table(
        rows, cols, left, top, width, row_h * rows,
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 배경색
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = t["table_header"]
            else:
                cell.fill.fore_color.rgb = (
                    t["table_row1"] if r % 2 == 1 else t["table_row2"]
                )

            # 텍스트 스타일
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(_cf)
                p.font.name = FONT_NAME
                p.font.bold = r == 0
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else t["text"]
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

    return table_shape


# ── 타이틀/섹션 ──────────────────────────────────────────────

def add_title_bar(
    slide, title: str, subtitle: str = "",
    bg_color: RGBColor | None = None,
    text_color: RGBColor | None = None,
    bar_height: float = 0.7,
    title_fontsize: int | None = None,
    subtitle_fontsize: int | None = None,
    theme: str | None = None,
) -> None:
    """상단 타이틀 바. theme='light' 또는 'dark'로 색상 자동 설정."""
    if theme == "light":
        bg = bg_color or THEME_LIGHT["primary"]
        text_color = text_color or RGBColor(0xFF, 0xFF, 0xFF)
    elif theme == "dark":
        bg = bg_color or THEME_DARK["primary"]
        text_color = text_color or RGBColor(0xFF, 0xFF, 0xFF)
    else:
        bg = bg_color or THEME_DARK["primary"]
        text_color = text_color or RGBColor(0xFF, 0xFF, 0xFF)
    _tf = title_fontsize or FONT_SIZES["title_bar"]
    _sf = subtitle_fontsize or FONT_SIZES["title_bar_sub"]
    _add_rect(
        slide, Inches(0), Inches(0),
        Inches(13.333), Inches(bar_height), bg,
    )
    add_textbox(
        slide, Inches(0.5), Inches(0.1),
        Inches(10), Inches(0.35),
        title, font_size=_tf, bold=True, color=text_color,
    )
    if subtitle:
        add_textbox(
            slide, Inches(0.5), Inches(0.4),
            Inches(10), Inches(0.25),
            subtitle, font_size=_sf, color=text_color,
        )


def add_section_divider(
    slide, left, top, width, title: str,
    color: RGBColor | None = None,
    title_fontsize: int | None = None,
) -> None:
    """섹션 구분선 (얇은 바 + 텍스트)."""
    c = color or THEME_DARK["primary"]
    _tf = title_fontsize or FONT_SIZES["section_divider"]
    _add_rect(slide, left, top, width, Pt(3), c)
    add_textbox(
        slide, left, top + Pt(6), width, Inches(0.3),
        title, font_size=_tf, bold=True, color=c,
    )


# ── 인사이트 박스 ────────────────────────────────────────────

def add_bullet_box(
    slide, left, top, width, height,
    items: list[str], title: str = "",
    positive: bool = True,
    card_bg: RGBColor | None = None,
    title_fontsize: int | None = None,
    item_fontsize: int | None = None,
) -> Any:
    """인사이트 박스 (Positive/Watch).

    Args:
        items: 불릿 항목 리스트
        positive: True면 초록, False면 빨강 액센트
    """
    accent = COLOR_POSITIVE if positive else COLOR_NEGATIVE
    bg = card_bg or RGBColor(0xF8, 0xF9, 0xFA)
    _tf = title_fontsize or FONT_SIZES["bullet_title"]
    _if = item_fontsize or FONT_SIZES["bullet_item"]
    card = add_rounded_rect(slide, left, top, width, height, bg)

    # 좌측 액센트 바
    _add_rect(slide, left, top, Pt(4), height, accent)

    y_offset = Inches(0.1)
    if title:
        add_textbox(
            slide, left + Inches(0.2), top + y_offset,
            width - Inches(0.4), Inches(0.25),
            title, font_size=_tf, bold=True, color=accent,
        )
        y_offset += Inches(0.3)

    # 불릿 항목
    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + y_offset,
        width - Inches(0.4), height - y_offset - Inches(0.1),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(_if)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.font.name = FONT_NAME
        p.space_after = Pt(4)

    return card


# ── 표지 슬라이드 ────────────────────────────────────────────

def add_cover_slide(
    prs: Presentation,
    title: str, subtitle: str, date_text: str,
    theme: str = "dark",
) -> Any:
    """표지 슬라이드 추가."""
    t = THEME_DARK if theme == "dark" else THEME_LIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, t["bg"])

    add_textbox(
        slide, Inches(1), Inches(2.5),
        Inches(11), Inches(1),
        title, font_size=FONT_SIZES["cover_title"], bold=True, color=t["text"],
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, Inches(1), Inches(3.8),
        Inches(11), Inches(0.5),
        subtitle, font_size=FONT_SIZES["cover_subtitle"], color=t["text_sub"],
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(1), Inches(4.5),
        Inches(11), Inches(0.4),
        date_text, font_size=FONT_SIZES["cover_date"], color=t["text_sub"],
        alignment=PP_ALIGN.CENTER,
    )
    return slide


# ── 차트 이미지 삽입 ─────────────────────────────────────────

def add_chart_image(
    slide, fig_or_buf, left, top, width, height,
) -> Any:
    """차트 이미지(Figure 또는 BytesIO)를 슬라이드에 삽입."""
    if hasattr(fig_or_buf, "savefig"):
        from src.util.chart_utils import chart_to_image
        buf = chart_to_image(fig_or_buf)
    else:
        buf = fig_or_buf
    return slide.shapes.add_picture(buf, left, top, width, height)
